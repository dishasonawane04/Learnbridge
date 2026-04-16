import fitz  # PyMuPDF
from pptx import Presentation
from PIL import Image
import os
import base64
import logging
from django.conf import settings

logger = logging.getLogger(__name__)

# Minimum characters from normal extraction before OCR fallback kicks in
_MIN_TEXT_CHARS = 100

_OCR_ERROR_MSG = (
    "Unable to detect text clearly. "
    "Please upload a clearer image or scan."
)


def extract_text_from_pdf(file_path):
    """
    Extract text from PDF using PyMuPDF.
    If the result is too short (scanned/handwritten PDF), fall back to OCR.
    """
    text = ""
    with fitz.open(file_path) as doc:
        for page in doc:
            text += page.get_text()

    # --- OCR Fallback for scanned / handwritten PDFs ---
    if len(text.strip()) < _MIN_TEXT_CHARS:
        logger.info(f"[OCR] PDF text too short ({len(text.strip())} chars). "
                    f"Attempting OCR fallback: {file_path}")
        from .ocr_utils import ocr_pdf_pages
        ocr_text = ocr_pdf_pages(file_path)
        if ocr_text and len(ocr_text.strip()) >= _MIN_TEXT_CHARS:
            logger.info(f"[OCR] Fallback succeeded: {len(ocr_text)} chars extracted.")
            return ocr_text
        else:
            logger.warning(f"[OCR] Fallback yielded insufficient text for: {file_path}")
            return _OCR_ERROR_MSG if not text.strip() else text

    return text


def extract_text_from_ppt(file_path):
    text = ""
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def extract_text_from_image(file_path):
    """
    Uses fast pytesseract for OCR first.
    Falls back to Ollama Vision model if PyTesseract fails to find anything.
    """
    import urllib.request
    import json
    from .ocr_utils import ocr_image

    logger.info(f"[OCR] Attempting pytesseract OCR: {file_path}")
    ocr_text = ocr_image(file_path)
    
    if ocr_text and len(ocr_text.strip()) >= 10:
        logger.info(f"[OCR] pytesseract succeeded: {len(ocr_text)} chars.")
        return ocr_text

    logger.info(f"[OCR] pytesseract returned little/no text. Attempting Ollama Vision fallback: {file_path}")
    ollama_text = ""
    try:
        with open(file_path, "rb") as image_file:
            encoded_string = base64.b64encode(image_file.read()).decode('utf-8')

        url = f"{settings.OLLAMA_BASE_URL}/api/generate"
        data = json.dumps({
            "model": getattr(settings, "OLLAMA_MODEL_VISION", "llava:latest"),
            "prompt": "Transcribe all text from this image accurately. Only return the transcribed text.",
            "images": [encoded_string],
            "stream": False,
            "options": {"num_predict": 500}
        }).encode('utf-8')

        req = urllib.request.Request(url, data=data, headers={'Content-Type': 'application/json'})
        logger.info(f"[OCR] Sending image to Ollama Vision (45s timeout): {file_path}")
        with urllib.request.urlopen(req, timeout=45) as response:
            result = json.loads(response.read().decode('utf-8'))
            ollama_text = result.get('response', '').strip()
    except Exception as e:
        logger.warning(f"[OCR] Ollama Vision failed for {file_path}: {e}")

    if ollama_text and len(ollama_text.strip()) >= 10:
        return ollama_text

    # Return whatever we got from pytesseract if it was short, or the error message
    if ocr_text:
        return ocr_text
        
    logger.warning(f"[OCR] Both pytesseract and Ollama failed for: {file_path}")
    return _OCR_ERROR_MSG

def extract_text_from_path(file_path):
    """Generic extraction based on file extension"""
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext in ['.pptx', '.ppt']:
        return extract_text_from_ppt(file_path)
    elif ext in ['.png', '.jpg', '.jpeg']:
        return extract_text_from_image(file_path)
    elif ext in ['.txt', '.md']:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    else:
        return ""

def extract_content(material):
    """Main entry point to extract text based on material type"""
    text = extract_text_from_path(material.file.path)
    
    # Set file type based on extension (legacy logic)
    ext = os.path.splitext(material.file.path)[1].lower()
    if ext == '.pdf': material.file_type = 'pdf'
    elif ext in ['.pptx', '.ppt']: material.file_type = 'ppt'
    elif ext in ['.png', '.jpg', '.jpeg']: material.file_type = 'image'
    elif ext in ['.txt', '.md']: material.file_type = 'text'
        
    material.extracted_text = text
    material.save()
    return text
